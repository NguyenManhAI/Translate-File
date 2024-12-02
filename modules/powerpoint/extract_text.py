from pptx import Presentation
import openai
from pptx.dml.color import _NoneColor

def extract_text_with_powerpoint(pptx_file,translate, src_lang, trg_lang):
    prs = Presentation(pptx_file)
    slides_data = []

    for slide in prs.slides:
        slide_info = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        font_color = run.font.color
                        font_color_value = None
                        if font_color:
                            print(font_color)
                            if isinstance(font_color, _NoneColor):
                                font_color_value = font_color.rgb
                                
                        run_info = {
                            "text": translate(run.text, src_lang, trg_lang) if run.text else None,
                            "bold": run.font.bold,
                            "italic": run.font.italic,
                            "font_size": run.font.size,
                            "font_name": run.font.name,
                            "font_color": font_color_value,
                        }
                        slide_info.append(run_info)
        slides_data.append(slide_info)

    return slides_data

