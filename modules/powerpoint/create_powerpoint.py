from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx import Presentation

# Bước 3: Tạo lại PowerPoint từ văn bản đã dịch và giữ nguyên định dạng
def create_translated_pptx(original_pptx_path, translated_data, output_pptx_path):
    prs = Presentation(original_pptx_path)
    
    slide_idx = 0
    for slide in prs.slides:
        text_idx = 0
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if text_idx < len(translated_data[slide_idx]):
                            translated_text_info = translated_data[slide_idx][text_idx]
                            
                            # Thay thế văn bản
                            run.text = translated_text_info["text"]
                            
                            # Khôi phục định dạng đã lưu
                            run.font.bold = translated_text_info["bold"]
                            run.font.italic = translated_text_info["italic"]
                            run.font.size = translated_text_info["font_size"]
                            run.font.name = translated_text_info["font_name"]
                            if translated_text_info["font_color"]:
                                run.font.color.rgb = translated_text_info["font_color"]
                            
                            text_idx += 1
        slide_idx += 1

    # Lưu file PowerPoint mới
    prs.save(output_pptx_path)